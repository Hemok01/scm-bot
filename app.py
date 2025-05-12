from flask import Flask, render_template, request, send_file
from pptx import Presentation
from io import BytesIO

app = Flask(__name__)


# 고정 헌금 계좌
DONATION_NOTE  = "<온라인 헌금> 우리 1005-804-217963 중앙대기독학생연합회"


# 고정 찬양 가사 및 제목
BLESSING_TITLE = "축복합니다"
BLESSING_LYRICS = """
축복합니다 주님의 이름으로
축복합니다 주님의 사랑으로

이 곳에 모인 주의 거룩한 자녀에게
주님의 기쁨과 주님의 사랑이
충만하게 충만하게 넘치기를 (축복합니다)

God bless you God bless you
축복합니다 주님의 사랑으로
"""

OFFERING_TITLE = "나의 모습 나의 소유"
OFFERING_LYRICS = """
나의 모습 나의 소유 주님 앞에 모두 드립니다
모든 아픔 모든 기쁨 내 모든 눈물 받아주소서

나의 생명을 드리니 
주 영광 위하여 사용하옵소서

내가 사는 날 동안에 
주를 찬양하며 기쁨의 제물되리
나를 받아주소서
"""

SENDING_TITLE = "가서 제자 삼으라"
SENDING_LYRICS = """
소나무 금잔디 동산에서
주님 젊은 제자들 다시 부르시사
마지막 그들에게 부탁하시기를
너희들은 가라 저 캠퍼스로

가서 제자삼으라 세상 많은 사람들을
세상 모든 영혼이 네게 달렸나니
가서 제자 삼으라 나의 길을 가르치라
내가 너희와 항상 함께 하리라
"""

def split_lyrics(lyrics: str, max_lines: int = 2):
    lines = [line for line in lyrics.splitlines() if line.strip()]
    return ["\n".join(lines[i:i + max_lines]) for i in range(0, len(lines), max_lines)]

def split_lines(text: str):
    return [line.strip() for line in text.splitlines() if line.strip()]

def generate_ppt(data: dict) -> BytesIO:
    prs = Presentation("template.pptx")

    # 1. 찬양과 경배 - 여러 곡
    titles = data.getlist("song_titles[]")
    lyrics = data.getlist("song_lyrics[]")
    for title, body in zip(titles, lyrics):
        slide_title = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title.placeholders[0].text = title
        for chunk in split_lyrics(body):
            slide_lyrics = prs.slides.add_slide(prs.slide_layouts[2])
            slide_lyrics.placeholders[0].text = chunk

    # 2. 기도
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.placeholders[13].text = "기도"
    slide.placeholders[12].text = data['대표 기도자']

    # 3. 성경봉독 (본문 + 구절 여러 슬라이드)
    verse_slide = prs.slides.add_slide(prs.slide_layouts[4])
    verse_slide.placeholders[13].text = "성경봉독"
    verse_slide.placeholders[12].text = data['성경 본문']

    for line in split_lines(data['성경 구절']):
        verse_text_slide = prs.slides.add_slide(prs.slide_layouts[3])
        verse_text_slide.placeholders[0].text = line

    # 4. 말씀
    sermon_slide = prs.slides.add_slide(prs.slide_layouts[4])
    sermon_slide.placeholders[13].text = data['설교 제목']
    sermon_slide.placeholders[12].text = data['설교자']

    # 5. 헌금 찬양
    slide_title = prs.slides.add_slide(prs.slide_layouts[1])
    slide_title.placeholders[0].text = OFFERING_TITLE
    for chunk in split_lyrics(OFFERING_LYRICS):
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        tf = slide.placeholders[0].text_frame
        tf.clear()

    # 6. 헌금기도
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.placeholders[13].text = "헌금기도"
    slide.placeholders[12].text = data['헌금 기도자']

    # 7. 환영
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.placeholders[13].text = "환영"
    slide.placeholders[12].text = data['환영 인도자']

    # 8. 축복송
    slide_title = prs.slides.add_slide(prs.slide_layouts[1])
    slide_title.placeholders[0].text = BLESSING_TITLE
    for chunk in split_lyrics(BLESSING_LYRICS):
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        slide.placeholders[0].text = chunk

    # 9. 광고
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.placeholders[13].text = "광고"
    slide.placeholders[12].text = data['광고자']

    # 10. 파송 찬양
    slide_title = prs.slides.add_slide(prs.slide_layouts[1])
    slide_title.placeholders[0].text = SENDING_TITLE
    for chunk in split_lyrics(SENDING_LYRICS):
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        slide.placeholders[0].text = chunk

    # 11. 축도
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    slide.placeholders[13].text = "축도"
    slide.placeholders[12].text = data['설교자']

    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        data = request.form
        ppt_io = generate_ppt(data)
        filename = data.get("filename", "Thursday_Service") + ".pptx"
        return send_file(ppt_io, as_attachment=True, download_name=filename)
    return render_template("form.html")
