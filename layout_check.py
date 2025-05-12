from pptx import Presentation

# --- 0. Layout inspection 유틸리티 추가 ---
def list_layout_placeholders(template_path: str = 'template.pptx'):
    """
    템플릿의 각 슬라이드 레이아웃에 정의된 플레이스홀더 목록을 출력합니다.
    """
    prs = Presentation(template_path)
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"Layout {idx}: '{layout.name}'")
        for shape in layout.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                print(f"  - Placeholder idx={ph.idx}, type={ph.type}, name='{shape.name}'")
        print()

list_layout_placeholders()
"""

# 프레젠테이션 로드
prs = Presentation('template.pptx')

def rename_placeholder_in_layout(prs, layout_idx: int, ph_idx: int, new_name: str):
    layout = prs.slide_layouts[layout_idx]
    for shape in layout.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == ph_idx:
            print(f"Renaming layout {layout_idx} placeholder {ph_idx} from '{shape.name}' to '{new_name}'")
            shape.name = new_name

# 예: 레이아웃 1(제목 및 콘텐츠)의 플레이스홀더 idx=1(본문 상자) 이름 변경

rename_placeholder_in_layout(prs, layout_idx=5, ph_idx=14, new_name="설교 제목")

# Layout 0: '제목 슬라이드'
rename_placeholder_in_layout(prs, layout_idx=0, ph_idx=0, new_name="제목 슬라이드 제목")
rename_placeholder_in_layout(prs, layout_idx=0, ph_idx=1, new_name="부제 슬라이드 부제목")
rename_placeholder_in_layout(prs, layout_idx=0, ph_idx=10, new_name="발행 날짜")
rename_placeholder_in_layout(prs, layout_idx=0, ph_idx=11, new_name="바닥글")
rename_placeholder_in_layout(prs, layout_idx=0, ph_idx=12, new_name="슬라이드 번호")

# Layout 1: '찬양 제목'
rename_placeholder_in_layout(prs, layout_idx=1, ph_idx=0, new_name="찬양 제목")

# Layout 2: '찬양 가사'
rename_placeholder_in_layout(prs, layout_idx=2, ph_idx=0, new_name="찬양 가사")

# Layout 3: '성경 구절'
rename_placeholder_in_layout(prs, layout_idx=3, ph_idx=0, new_name="성경 구절")

# Layout 4: '설교말씀 오프닝'
rename_placeholder_in_layout(prs, layout_idx=4, ph_idx=12, new_name="설교 말씀")
rename_placeholder_in_layout(prs, layout_idx=4, ph_idx=13, new_name="설교 제목")

# Layout 5: '설교제목, 목사님 성함 우측에'
rename_placeholder_in_layout(prs, layout_idx=5, ph_idx=14, new_name="설교 제목")
rename_placeholder_in_layout(prs, layout_idx=5, ph_idx=13, new_name="설교자 이름")

# Layout 6: '양장님 환영'
rename_placeholder_in_layout(prs, layout_idx=6, ph_idx=12, new_name="환영")
rename_placeholder_in_layout(prs, layout_idx=6, ph_idx=13, new_name="양장 이름(학과 학번)")

# Layout 7: '새가족 환영'
rename_placeholder_in_layout(prs, layout_idx=7, ph_idx=12, new_name="새가족 환영")
rename_placeholder_in_layout(prs, layout_idx=7, ph_idx=13, new_name="새가족 이름(학과 학번)")

# Layout 8: '헌금기도'
rename_placeholder_in_layout(prs, layout_idx=8, ph_idx=12, new_name="헌금기도")
rename_placeholder_in_layout(prs, layout_idx=8, ph_idx=13, new_name="헌금기도자")

# Layout 9: '광고'
rename_placeholder_in_layout(prs, layout_idx=9, ph_idx=12, new_name="광고")
rename_placeholder_in_layout(prs, layout_idx=9, ph_idx=13, new_name="관장 이름(학과 학번)")

# Layout 10: '축도'
rename_placeholder_in_layout(prs, layout_idx=10, ph_idx=12, new_name="축도")
rename_placeholder_in_layout(prs, layout_idx=10, ph_idx=13, new_name="설교자 이름")

# Layout 11: '제목 및 내용'
# No placeholders to rename or use default names

# 저장
prs.save('template.pptx')

if __name__ == "__main__":
    # --- GUI Setup ---
    root = tk.Tk()
    root.title("목요예배 PPT 생성기")

    # Song Title
    tk.Label(root, text="찬양 제목:").grid(row=0, column=0, sticky="e")
    entry_title = tk.Entry(root, width=40)
    entry_title.grid(row=0, column=1, padx=5, pady=5)

    # Lyrics
    tk.Label(root, text="찬양 가사:").grid(row=1, column=0, sticky="ne")
    text_lyrics = scrolledtext.ScrolledText(root, width=40, height=10)
    text_lyrics.grid(row=1, column=1, padx=5, pady=5)

    # Service Info
    info_labels = ["찬양 인도자", "대표 기도자", "설교자", "주제 말씀", "설교 제목"]
    entries_info = {}
    for i, label in enumerate(info_labels, start=2):
        tk.Label(root, text=label + ":").grid(row=i, column=0, sticky="e")
        ent = tk.Entry(root, width=40)
        ent.grid(row=i, column=1, padx=5, pady=2)
        entries_info[label] = ent

    def on_generate():
        song_title = entry_title.get().strip()
        lyrics = text_lyrics.get("1.0", tk.END).strip()
        if not song_title or not lyrics:
            messagebox.showwarning("입력 오류", "찬양 제목과 가사를 모두 입력하세요.")
            return
        chunks = split_lyrics(lyrics)
        service_info = {label: entries_info[label].get().strip() for label in info_labels}
        try:
            create_service_ppt(song_title, chunks, service_info)
            messagebox.showinfo("완료", "PPT 파일이 생성되었습니다.")
        except Exception as e:
            messagebox.showerror("오류", str(e))

    btn = tk.Button(root, text="PPT 생성", command=on_generate)
    btn.grid(row=7, column=1, pady=10)

    root.mainloop()
    """