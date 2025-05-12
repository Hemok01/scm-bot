import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- 1. 가사 입력 및 분할 ---
def split_lyrics(lyrics: str, max_lines_per_slide: int = 6) -> list[str]:
    """
    Split lyrics by lines and group into slides.
    """
    lines = [line for line in lyrics.splitlines() if line.strip()]
    slides = []
    for i in range(0, len(lines), max_lines_per_slide):
        slides.append("\n".join(lines[i:i + max_lines_per_slide]))
    return slides

# --- 2. 예배 정보 입력 ---
def collect_service_info() -> dict:
    """
    Prompt user for worship service details.
    """
    info = {}
    info['찬양 제목'] = input("찬양 제목: ", default="찬양 제목")
    info['찬양 인도자'] = input("찬양 인도자: ", default="찬양 인도자")
    info['대표 기도자'] = input("대표 기도자: ", default="대표 기도자")
    info['설교자'] = input("설교자: ", default="설교자")
    info['주제 말씀'] = input("주제 말씀: ", default="주제 말씀")
    info['성경 구절'] = input("성경 구절: ", default="성경 구절")
    info['설교 제목'] = input("설교 제목: ", default="설교 제목")
    info['헌금 기도자'] = input("헌금 기도자: ", default="헌금 기도자")
    info['관장 이름'] = input("관장 이름(학과 학번): ", default="관장 이름")
    info['양장 이름'] = input("양장 이름(학과 학번): ", default="양장 이름")
    info['헌금 기도자'] = input("헌금 기도자: ", default="헌금 기도자")
    info['새가족 이름'] = input("새가족 이름(학과 학번): ", default="새가족 이름")
    return info

# --- 3. PPT 생성 ---
def create_service_ppt(lyrics_chunks: list[str], service_info: dict, output_path: str = "Thursday_Service.pptx"):
    # load existing template with pre‑configured slide master
    prs = Presentation('template.pptx')

    # 3-2. 찬양 제목 슬라이드
    slide_title = prs.slides.add_slide(prs.slide_layouts[1])
    slide_title.placeholders[0].text = service_info['찬양 제목']

    # 3-3. 찬양 가사 슬라이드
    for chunk in lyrics_chunks:
        slide_lyrics = prs.slides.add_slide(prs.slide_layouts[2])
        slide_lyrics.placeholders[0].text = chunk

    # 3-4. 성경 봉독 슬라이드
    slide_verse = prs.slides.add_slide(prs.slide_layouts[4])
    slide_verse.placeholders[13].text = "-성경 봉독-"
    slide_verse.placeholders[12].text = service_info['주제 말씀']

    # 3-5. 성경 구절 슬라이드
    slide_verse_text = prs.slides.add_slide(prs.slide_layouts[3])
    slide_verse_text.placeholders[0].text = service_info['성경 구절']

    # 3-1. 제목 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[12].text = "-설교 말씀-"
    slide.placeholders[13].text = service_info['설교 제목']

    # 3-4. 헌금기도 슬라이드
    slide_offering = prs.slides.add_slide(prs.slide_layouts[8])
    slide_offering.placeholders[13].text = "-헌금 기도-"
    slide_offering.placeholders[12].text = service_info['헌금 기도자']

    # 3-6. 환영 슬라이드
    slide_welcome = prs.slides.add_slide(prs.slide_layouts[6])
    slide_welcome.placeholders[13].text = "-환영-"
    slide_welcome.placeholders[12].text = service_info['양장 이름(학과 학번)']

    # 3-7. 새가족 환영 슬라이드
    slide_newcomer = prs.slides.add_slide(prs.slide_layouts[7])
    slide_newcomer.placeholders[13].text = "-새가족 환영-"
    slide_newcomer.placeholders[12].text = service_info['새가족 이름(학과 학번)']

    # 3-5. 광고 슬라이드
    slide_notice = prs.slides.add_slide(prs.slide_layouts[9])
    slide_notice.placeholders[13].text = "-광고-"
    slide_notice.placeholders[12].text = service_info['관장 이름']

    # 3-6. 축도 슬라이드
    slide_benediction = prs.slides.add_slide(prs.slide_layouts[10])
    slide_benediction.placeholders[13].text = "-축도-"
    slide_benediction.placeholders[12].text = service_info['설교자']

    prs.save(output_path)
    print(f"Saved PPT: {output_path}")

# --- 4. Main 실행 ---
if __name__ == "__main__":
    # 가사 입력
    raw = input("가사를 붙여넣으세요: \n")
    chunks = split_lyrics(raw)

    # 예배 정보 입력
    info = collect_service_info()

    # PPT 생성
    create_service_ppt(chunks, info)
